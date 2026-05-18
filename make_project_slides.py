from __future__ import annotations

import argparse
import json
from pathlib import Path


def _fmt_pct(x: float) -> str:
    return f"{x*100:.1f}%"


def _load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def build_slides(*, report_dir: Path, out_pptx: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    report_dir = report_dir.resolve()
    metrics_path = report_dir / "metrics.json"
    hfa_path = report_dir / "home_field_effect.json"
    model_metrics_path = Path("models") / "metrics.json"

    metrics = _load_json(metrics_path) if metrics_path.exists() else {}
    hfa = _load_json(hfa_path) if hfa_path.exists() else {}
    model_metrics = _load_json(model_metrics_path) if model_metrics_path.exists() else {}

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # Theme-ish colors
    GT_NAVY = RGBColor(0x00, 0x30, 0x57)
    GT_GOLD = RGBColor(0xB3, 0xA3, 0x69)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED = RGBColor(0xDD, 0xDD, 0xDD)

    def add_title(slide, title: str, subtitle: str | None = None) -> None:
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(1.1))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = WHITE
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = MUTED

    def add_bullets(slide, x, y, w, h, title: str, bullets: list[str]) -> None:
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text = title
        t.font.size = Pt(20)
        t.font.bold = True
        t.font.color.rgb = GT_GOLD
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(15)
            p.font.color.rgb = WHITE

    def add_small_bullets(slide, x, y, w, h, title: str, bullets: list[str]) -> None:
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.clear()
        t = tf.paragraphs[0]
        t.text = title
        t.font.size = Pt(18)
        t.font.bold = True
        t.font.color.rgb = GT_GOLD
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE

    def add_footer(slide, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(12.2), Inches(0.35))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.RIGHT

    def set_bg(slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = GT_NAVY

    # Pull core numbers used across multiple slides
    rec = (metrics.get("record") or {}) if isinstance(metrics, dict) else {}
    marg = (metrics.get("margin") or {}) if isinstance(metrics, dict) else {}
    splits = (metrics.get("splits") or {}) if isinstance(metrics, dict) else {}

    wins = rec.get("wins")
    losses = rec.get("losses")
    win_rate = rec.get("win_rate")
    win_ci = (rec.get("win_rate_wilson_95") or {}) if isinstance(rec.get("win_rate_wilson_95"), dict) else {}
    mean_margin = marg.get("mean")
    margin_ci = (marg.get("mean_bootstrap_95") or {}) if isinstance(marg.get("mean_bootstrap_95"), dict) else {}

    # Model metrics (train/val/test) if available
    train_cfg = (model_metrics.get("train_config") or {}) if isinstance(model_metrics, dict) else {}
    val_cal = (((model_metrics.get("val") or {}).get("calibrated")) or {}) if isinstance(model_metrics, dict) else {}
    test_cal = (((model_metrics.get("test") or {}).get("calibrated")) or {}) if isinstance(model_metrics, dict) else {}

    # Home-field estimate
    hfa_margin = (hfa.get("margin_model") or {}) if isinstance(hfa, dict) else {}
    hfa_win = (hfa.get("win_model") or {}) if isinstance(hfa, dict) else {}

    # Static image paths
    img_winrate = report_dir / "winrate_by_year.png"
    img_margin_by_year = report_dir / "margin_by_year.png"
    img_margin_hist = report_dir / "margin_hist.png"
    img_close = report_dir / "close_games_by_year.png"
    img_blowouts = report_dir / "blowouts_by_year.png"
    img_pythag = report_dir / "pythag_vs_actual.png"
    img_sos = report_dir / "opp_strength_by_year.png"
    img_elo_gap = report_dir / "elo_gap_vs_win.png"

    # ---- Slide 1: Title ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    add_title(
        slide,
        "GT Win Probability Project (10-slide overview)",
        "End-to-end: data engineering → baseline ML model → fan-friendly dashboard (2014–2025)",
    )
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.7),
        Inches(12.0),
        Inches(4.9),
        "What we built (in one pipeline)",
        [
            "A repeatable pipeline that pulls CFBD data → builds a clean dataset → trains a model → produces a dashboard",
            "Goal: predict Georgia Tech win probability without using CFBD’s win expectancy",
            "Deliverables: dataset CSV, trained model artifacts, season predictions CSVs, dashboard HTML + charts",
        ],
    )
    add_footer(slide, "Files: winprob.py • analysis_report.py • reports/gt_summary_2014_2025/dashboard.html")

    # ---- Slide 2: Problem statement & scope ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Slide 2/10: Problem + Scope", "What question are we answering?")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(6.1),
        Inches(5.4),
        "Core question",
        [
            "For a given matchup, what is GT’s pregame chance to win? (a probability, not just W/L)",
            "We use only information available before kickoff (no postgame stats as inputs)",
            "We do not use betting spread as a feature",
        ],
    )
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.55),
        Inches(5.6),
        Inches(5.4),
        "What’s included",
        [
            "Historical games (2014–2025 regular season) to build a training dataset",
            "Baseline model + calibration (probabilities should be meaningful)",
            "A dashboard to explain results to fans / teammates",
        ],
    )
    add_footer(slide, "Why probability? It supports decisions, comparisons, and “how surprised should we be?”")

    # ---- Slide 3: Data sources ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Slide 3/10: Data Sources (CFBD API)", "Where the raw data comes from")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(6.1),
        Inches(5.3),
        "Main endpoints used",
        [
            "`/games`: schedule, opponent, location, scores (for label)",
            "`/ratings/elo`: team strength proxy (pregame-ish)",
            "`/talent`, `/recruiting/teams`, `/player/returning`: roster strength proxies",
            "`/teams`: team name resolution + FBS/FCS classification",
            "Optional: `/lines` for market-implied probabilities (benchmarking, not training)",
        ],
    )
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.5),
        Inches(5.6),
        Inches(5.3),
        "Why multiple sources",
        [
            "A single endpoint rarely has everything you need for good pregame features",
            "We combine strength + roster + season-to-date form to get a better signal",
            "We cache responses to keep runs fast and reproducible",
        ],
    )
    add_footer(slide, "Key component: CFBDClient caches JSON in data_raw/cfbd_cache/")

    # ---- Slide 4: Data engineering pipeline (how we turn JSON into a dataset) ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Slide 4/10: Data Engineering Pipeline", "Raw API JSON → clean ML-ready table")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(6.1),
        Inches(5.4),
        "Pipeline steps",
        [
            "Call CFBD endpoints (with caching) and resolve team names consistently",
            "Build one row per GT game: (year, week, opponent, home/away/neutral)",
            "Compute pregame features using only games before that week (season-to-date form)",
            "Write the dataset to CSV so training is decoupled from API calls",
        ],
    )
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.55),
        Inches(5.6),
        Inches(5.4),
        "Important engineering details",
        [
            "Label is `gt_win` (computed from final score) — used only as the target",
            "Non-FBS opponents can have missing Elo/talent → we apply a floor to avoid bad imputations",
            "Missing values are expected; we impute numerics during training (median)",
            "Outputs: data_processed/model_dataset.csv + cached raw JSON",
        ],
    )
    add_footer(slide, "Command: python winprob.py build-dataset --from-year 2014 --to-year 2025")

    # ---- Slide 5: Dataset schema & features ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Slide 5/10: Dataset + Feature Engineering", "What the model actually sees before a game")
    add_small_bullets(
        slide,
        Inches(0.8),
        Inches(1.6),
        Inches(6.2),
        Inches(5.2),
        "Key columns (examples)",
        [
            "Identifiers: game_id, year, week, opponent, location flags",
            "Target label: gt_win (1/0)",
            "Season-to-date form diffs: win_pct_diff, point_diff_pg_diff, points_for/against diffs",
            "Weighted “recent form” diffs: w_* versions (recent games count more)",
            "Strength/roster diffs: elo_diff, talent_diff, returning_diff, recruiting diffs",
            "Opponent type: opp_is_fbs (handles FCS games more realistically)",
        ],
    )
    # Visual: margin histogram helps explain why W/L alone is noisy
    if img_margin_hist.exists():
        slide.shapes.add_picture(str(img_margin_hist), Inches(7.1), Inches(1.75), width=Inches(5.8))
    add_footer(slide, "Pregame-only rule: features use weeks < current week, not postgame stats")

    # ---- Slide 6: Model choice & training setup ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Slide 6/10: Modeling Approach", "Simple baseline: logistic regression + calibration")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.6),
        Inches(6.1),
        Inches(5.2),
        "How it works",
        [
            "Supervised learning: features → probability GT wins",
            "Train/val/test split by season (avoid mixing years): 2014–2022 / 2023–2024 / 2025",
            f"Recency weighting so recent years matter more (half-life ≈ {train_cfg.get('half_life_years', 'N/A')} years)",
            "Pipeline: impute missing values → standardize → logistic regression",
            f"Calibration: {train_cfg.get('calibration', 'auto')} (so 0.70 means ‘should win ~70% of the time’)",
        ],
    )
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.6),
        Inches(5.6),
        Inches(5.2),
        "Why this baseline is good",
        [
            "Explainable: you can tell a story about what drives probability (home field, strength gap, form)",
            "Fast to train and easy to debug (good for a first version)",
            "Sets a benchmark: future models should beat this to justify complexity",
        ],
    )
    add_footer(slide, "Train command: python winprob.py train  (writes models/metrics.json)")

    # ---- Slide 7: Model evaluation (what the metrics mean) ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Slide 7/10: Evaluation", "Did the probabilities make sense?")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(6.2),
        Inches(5.4),
        "Metrics you can explain",
        [
            "Accuracy: how often we picked the winner (threshold 50%)",
            "Brier score: how close probabilities were to outcomes (lower is better)",
            "Log loss: penalizes confident wrong predictions more than small mistakes",
            "Calibration: do predicted probabilities match real win rates over time?",
        ],
    )
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.55),
        Inches(5.6),
        Inches(5.4),
        "Latest run (from models/metrics.json)",
        [
            f"Validation (2023–2024) accuracy: {val_cal.get('accuracy', float('nan')):.3f}",
            f"Validation Brier: {val_cal.get('brier', float('nan')):.3f}",
            f"Test (2025) accuracy: {test_cal.get('accuracy', float('nan')):.3f}",
            f"Test Brier: {test_cal.get('brier', float('nan')):.3f}",
            "Caveat: test is ~12 games → treat as directional, not definitive",
        ],
    )
    add_footer(slide, "Best practice: keep 2025 as the clean out-of-sample test year")

    # ---- Slide 8: Dashboard overview (what it contains) ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Slide 8/10: Dashboard", "Turning results into something people can use")
    add_small_bullets(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(6.2),
        Inches(5.4),
        "What the dashboard answers",
        [
            f"Overall record: {wins}-{losses} (win rate {_fmt_pct(win_rate) if isinstance(win_rate, (int, float)) else 'N/A'})",
            "Was GT consistently good, or did it swing by year? (win rate by season)",
            "How ‘dominant’ were games? (margin distribution + margin by year)",
            "Does home field matter? (splits + regression estimate)",
            "Context: close games, blowouts, schedule strength proxy, Elo gap vs win rate",
        ],
    )
    # Right side: 2 visuals
    y_img = Inches(1.75)
    if img_winrate.exists():
        slide.shapes.add_picture(str(img_winrate), Inches(7.1), y_img, width=Inches(5.8))
        y_img += Inches(2.65)
    if img_margin_by_year.exists():
        slide.shapes.add_picture(str(img_margin_by_year), Inches(7.1), y_img, width=Inches(5.8))
    add_footer(slide, f"Open: {report_dir.name}/dashboard.html (static PNGs included for offline)")

    # ---- Slide 9: Example deeper insight (home-field) ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Slide 9/10: Example Insight", "Home-field effect (simple regression)")
    home_rate = (splits.get("home") or {}).get("rate") if isinstance(splits.get("home"), dict) else None
    away_rate = (splits.get("away") or {}).get("rate") if isinstance(splits.get("away"), dict) else None
    hfa_pts = hfa_margin.get("home_adv_points")
    hfa_pts_ci = (hfa_margin.get("home_adv_points_bootstrap_95") or {}) if isinstance(hfa_margin.get("home_adv_points_bootstrap_95"), dict) else {}
    hfa_shift = hfa_win.get("p_win_even_shift")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(6.2),
        Inches(5.4),
        "What it means",
        [
            f"Home win rate ≈ {_fmt_pct(home_rate)} vs away win rate ≈ {_fmt_pct(away_rate)}"
            if isinstance(home_rate, (int, float)) and isinstance(away_rate, (int, float))
            else "Home vs away win rates: see dashboard splits",
            f"Regression estimate: home field ≈ +{hfa_pts:.1f} points"
            if isinstance(hfa_pts, (int, float))
            else "Regression estimate: home field points boost (see report)",
            (
                f"Uncertainty (95% CI): [{hfa_pts_ci.get('low'):.1f}, {hfa_pts_ci.get('high'):.1f}] points"
                if isinstance(hfa_pts_ci.get("low"), (int, float)) and isinstance(hfa_pts_ci.get("high"), (int, float))
                else "Uncertainty: wide CI because samples are limited"
            ),
            f"Win-prob shift when teams are ‘even’: about +{_fmt_pct(hfa_shift)}"
            if isinstance(hfa_shift, (int, float))
            else "Win-prob shift: computed in home_field_effect.json",
            "Key takeaway: looks helpful, but the uncertainty means it’s suggestive, not a guarantee",
        ],
    )
    # Add one more chart for context
    if img_sos.exists():
        slide.shapes.add_picture(str(img_sos), Inches(7.1), Inches(1.85), width=Inches(5.8))
    add_footer(slide, "Why this slide matters: shows we can extract interpretable insights, not just a black-box model")

    # ---- Slide 10: Model vs market + challenges + roadmap ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Slide 10/10: Benchmarking + Next Steps", "How we judge it and what’s next")
    # Win-rate & margin context lines
    win_ci_line = ""
    if isinstance(win_ci.get("low"), (int, float)) and isinstance(win_ci.get("high"), (int, float)):
        win_ci_line = f"Overall win rate uncertainty (95% CI): {_fmt_pct(win_ci['low'])}–{_fmt_pct(win_ci['high'])}"
    margin_ci_line = ""
    if isinstance(margin_ci.get("low"), (int, float)) and isinstance(margin_ci.get("high"), (int, float)):
        margin_ci_line = f"Margin uncertainty (95% CI): {margin_ci['low']:+.2f} to {margin_ci['high']:+.2f} points/game"

    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(6.2),
        Inches(5.4),
        "How we benchmark + what can go wrong",
        [
            "Benchmark idea: compare to betting market implied win prob (strong public baseline)",
            "Small samples: a 12-game season makes “accuracy” swing a lot",
            "Non-stationary team: coaches/players change → older seasons matter less (we use recency weighting)",
            "Missing/uneven data (especially non-FBS opponents) can distort features",
            "Football randomness: close games + turnovers mean probabilities will never be perfect",
        ],
    )
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.55),
        Inches(5.6),
        Inches(5.4),
        "Next improvements (roadmap)",
        [
            "Train on all FBS games, then predict GT (massively larger dataset + better generalization)",
            "Add stronger pregame team strength features (and better schedule-adjusted form)",
            "Try non-linear models (gradient boosting) and tune with cross-validation",
            "Improve evaluation: reliability diagrams, rolling-year backtests, more seasons for test",
            "Product: add 2024 model-vs-market table + richer matchup explanations in the dashboard",
            *([""] if not win_ci_line and not margin_ci_line else []),
            *(["Context stats: " + win_ci_line] if win_ci_line else []),
            *(["Context stats: " + margin_ci_line] if margin_ci_line else []),
        ],
    )
    if img_elo_gap.exists():
        slide.shapes.add_picture(str(img_elo_gap), Inches(7.05), Inches(5.75), width=Inches(5.85))
    add_footer(slide, "Positioning: strong baseline pipeline + explainable dashboard; ready to extend")

    # ---- Slide 11: Appendix (optional) ----
    # Keep deck at exactly 10 slides by removing this slide if added later.

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))


def main() -> int:
    ap = argparse.ArgumentParser(description="Generate a PPTX slide deck summarizing the GT project.")
    ap.add_argument(
        "--report-dir",
        type=Path,
        default=Path("reports") / "gt_summary_2014_2025",
        help="Report directory containing dashboard assets (default: reports/gt_summary_2014_2025)",
    )
    ap.add_argument(
        "--out",
        type=Path,
        default=Path("reports") / "gt_summary_2014_2025" / "project_overview_slides.pptx",
        help="Output pptx path",
    )
    args = ap.parse_args()
    build_slides(report_dir=args.report_dir, out_pptx=args.out)
    print(f"Wrote slides to {args.out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
